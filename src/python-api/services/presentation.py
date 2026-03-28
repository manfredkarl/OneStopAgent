"""PowerPoint generation service — PptxGenJS (primary) + python-pptx (fallback)."""

import logging
import os
import subprocess
import tempfile
import uuid

logger = logging.getLogger(__name__)

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "output")
# node_modules lives next to this service's parent package.json
NODE_MODULES = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "node_modules"))


def execute_pptxgenjs(script: str, customer_name: str = "Customer") -> str:
    """Write a PptxGenJS script to a temp file, execute it, return the .pptx path.

    The script must use the literal string ``OUTPUT_PATH`` as the fileName
    argument to ``pres.writeFile()``. This function replaces it with the
    real output path before execution.
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    safe_name = "".join(c if c.isalnum() or c in "-_ " else "" for c in customer_name).strip().replace(" ", "_")
    filename = f"OneStopAgent-{safe_name}-{uuid.uuid4().hex[:8]}.pptx"
    output_path = os.path.join(OUTPUT_DIR, filename)

    # Inject the real output path — use forward slashes for Node on all platforms
    abs_output = os.path.abspath(output_path).replace("\\", "/")
    script = script.replace("OUTPUT_PATH", f'"{abs_output}"')

    # Write script to temp file
    script_fd, script_path = tempfile.mkstemp(suffix=".js", prefix="pptxgen_")
    try:
        with os.fdopen(script_fd, "w", encoding="utf-8") as f:
            f.write(script)

        env = {**os.environ, "NODE_PATH": NODE_MODULES}
        result = subprocess.run(
            ["node", script_path],
            capture_output=True,
            text=True,
            timeout=60,
            cwd=OUTPUT_DIR,
            env=env,
        )

        if result.returncode != 0:
            logger.error("PptxGenJS script failed (exit %d):\nstdout: %s\nstderr: %s",
                         result.returncode, result.stdout[:500], result.stderr[:500])
            raise RuntimeError(f"PptxGenJS script failed: {result.stderr[:300]}")

        if not os.path.isfile(output_path):
            raise FileNotFoundError(f"PptxGenJS ran but output not found at {output_path}")

        logger.info("Generated PPTX: %s", output_path)
        return output_path

    finally:
        # Clean up temp script
        try:
            os.unlink(script_path)
        except OSError:
            pass


# ---------------------------------------------------------------------------
# python-pptx fallback — used when LLM / Node.js path is unavailable
# ---------------------------------------------------------------------------

# Colors (RGBColor hex)
_DARK_BG = "1E2761"
_ACCENT = "0891B2"
_LIGHT_BG = "F8FAFC"
_TEXT = "1E293B"
_SUBTEXT = "64748B"
_TABLE_HDR = "1E2761"
_TABLE_BORDER = "E2E8F0"


def create_pptx_python(slide_data: dict, customer_name: str = "Customer") -> str:
    """Create a professional PPTX deck using python-pptx (no LLM / Node required).

    Returns the absolute path to the generated .pptx file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    safe_name = "".join(c if c.isalnum() or c in "-_ " else "" for c in customer_name).strip().replace(" ", "_")
    filename = f"OneStopAgent-{safe_name}-{uuid.uuid4().hex[:8]}.pptx"
    output_path = os.path.abspath(os.path.join(OUTPUT_DIR, filename))

    # Use branded template if available, otherwise blank
    TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "templates", "slide_master.pptx")
    if os.path.exists(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
        # Remove existing example slides from the template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId:
                prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        logger.info(f"Using branded template: {TEMPLATE_PATH}")
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        logger.info("No template found — using blank presentation")
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    # Use template layouts if available, fallback to index-based
    def _get_layout(name_hint: str, fallback_idx: int = 0):
        """Find layout by name, fallback to index."""
        for layout in prs.slide_layouts:
            if name_hint.lower() in layout.name.lower():
                return layout
        if fallback_idx < len(prs.slide_layouts):
            return prs.slide_layouts[fallback_idx]
        return prs.slide_layouts[0]

    title_layout = _get_layout("Title", 0)
    content_layout = _get_layout("Title and Content", 6)
    divider_layout = _get_layout("Divider", 3)
    blank_layout = _get_layout("Blank", 27)

    def _rgb(hex_str: str) -> RGBColor:
        return RGBColor.from_string(hex_str)

    def _add_bg(slide, color_hex: str):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(color_hex)

    def _add_textbox(slide, left, top, width, height, text, font_size=14,
                     bold=False, color=_TEXT, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = _rgb(color)
        p.font.name = font_name
        p.alignment = alignment
        return tf

    def _add_bullets(slide, left, top, width, height, items, font_size=13, color=_TEXT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(font_size)
            p.font.color.rgb = _rgb(color)
            p.font.name = "Calibri"
            p.level = 0
            p.space_after = Pt(6)
            pf = p._pPr
            if pf is None:
                from pptx.oxml.ns import qn
                pf = p._p.get_or_add_pPr()
            from pptx.oxml.ns import qn
            buNone = pf.findall(qn("a:buNone"))
            for el in buNone:
                pf.remove(el)
            from lxml import etree
            buChar = etree.SubElement(pf, qn("a:buChar"))
            buChar.set("char", "•")

    def _add_rect(slide, left, top, width, height, color_hex):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(color_hex)
        shape.line.fill.background()
        return shape

    def _add_stat_card(slide, left, top, width, height, value, label, bg_color="FFFFFF", value_color=_ACCENT):
        card = _add_rect(slide, left, top, width, height, bg_color)
        _add_textbox(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.8),
                     str(value), font_size=36, bold=True, color=value_color, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, left + Inches(0.2), top + Inches(1.0), width - Inches(0.4), Inches(0.5),
                     label, font_size=12, color=_SUBTEXT, alignment=PP_ALIGN.CENTER)

    customer = slide_data.get("customer", customer_name)
    problem = slide_data.get("problem", "")

    # ── Slide 1: Title ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, _DARK_BG)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
    _add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 f"Azure Solution Proposal", font_size=40, bold=True, color="FFFFFF", alignment=PP_ALIGN.LEFT)
    _add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
                 customer, font_size=28, bold=False, color=_ACCENT, alignment=PP_ALIGN.LEFT)
    _add_textbox(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.5),
                 "Generated by OneStopAgent", font_size=12, color=_SUBTEXT, alignment=PP_ALIGN.LEFT)

    # ── Slide 2: Problem & Opportunity ───────────────────────────────────
    if problem:
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, _LIGHT_BG)
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                     "Problem & Opportunity", font_size=30, bold=True, color=_DARK_BG)
        _add_rect(slide, Inches(0.8), Inches(1.3), Inches(0.08), Inches(1.8), _ACCENT)
        _add_textbox(slide, Inches(1.2), Inches(1.3), Inches(10), Inches(2.5),
                     problem, font_size=16, color=_TEXT)
        clarifications = slide_data.get("clarifications", "")
        if clarifications:
            _add_textbox(slide, Inches(1.2), Inches(4.0), Inches(10), Inches(2),
                         clarifications, font_size=13, color=_SUBTEXT)

    # ── Slide 3: Solution Overview ───────────────────────────────────────
    arch = slide_data.get("architecture")
    if arch:
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, _LIGHT_BG)
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                     "Solution Overview", font_size=30, bold=True, color=_DARK_BG)

        narrative = arch.get("narrative", "")
        based_on = arch.get("basedOn", "")
        body = narrative
        if based_on:
            body += f"\n\nBased on: {based_on}"
        _add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(2.5),
                     body, font_size=15, color=_TEXT)

    # ── Slide 4: Architecture Components ─────────────────────────────────
    components = arch.get("components", []) if arch else []
    if components:
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, _LIGHT_BG)
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                     "Architecture Components", font_size=30, bold=True, color=_DARK_BG)

        from pptx.util import Emu
        tbl_data = [["Component", "Azure Service", "Description"]]
        for c in components:
            tbl_data.append([
                c.get("name", ""),
                c.get("azureService", ""),
                c.get("description", ""),
            ])

        rows, cols = len(tbl_data), 3
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5 * rows))
        tbl = table_shape.table
        tbl.columns[0].width = Inches(2.5)
        tbl.columns[1].width = Inches(3)
        tbl.columns[2].width = Inches(6)
        for r_idx, row_data in enumerate(tbl_data):
            for c_idx, val in enumerate(row_data):
                cell = tbl.cell(r_idx, c_idx)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.name = "Calibri"
                    if r_idx == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = _rgb("FFFFFF")
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(_TABLE_HDR)
                    else:
                        paragraph.font.color.rgb = _rgb(_TEXT)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb("FFFFFF" if r_idx % 2 == 1 else "F1F5F9")

    # ── Slide 5: Azure Services & SKUs ───────────────────────────────────
    services = slide_data.get("services", [])
    if services:
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, _LIGHT_BG)
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                     "Azure Services & SKUs", font_size=30, bold=True, color=_DARK_BG)

        tbl_data = [["Service", "SKU", "Region"]]
        for s in services:
            tbl_data.append([s.get("name", ""), s.get("sku", ""), s.get("region", "")])

        rows, cols = len(tbl_data), 3
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5 * rows))
        tbl = table_shape.table
        tbl.columns[0].width = Inches(4)
        tbl.columns[1].width = Inches(3.5)
        tbl.columns[2].width = Inches(4)
        for r_idx, row_data in enumerate(tbl_data):
            for c_idx, val in enumerate(row_data):
                cell = tbl.cell(r_idx, c_idx)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.name = "Calibri"
                    if r_idx == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = _rgb("FFFFFF")
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(_TABLE_HDR)
                    else:
                        paragraph.font.color.rgb = _rgb(_TEXT)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb("FFFFFF" if r_idx % 2 == 1 else "F1F5F9")

    # ── Slide 6: Cost Estimate ───────────────────────────────────────────
    costs = slide_data.get("costs")
    if costs:
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, _LIGHT_BG)
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                     "Cost Estimate", font_size=30, bold=True, color=_DARK_BG)

        monthly = costs.get("totalMonthly", 0)
        annual = costs.get("totalAnnual", 0)
        _add_stat_card(slide, Inches(0.8), Inches(1.4), Inches(3.5), Inches(1.6),
                       f"${monthly:,.0f}/mo", "Monthly Cost")
        _add_stat_card(slide, Inches(4.8), Inches(1.4), Inches(3.5), Inches(1.6),
                       f"${annual:,.0f}/yr", "Annual Cost")

        pricing_src = costs.get("pricingSource", "")
        if pricing_src:
            _add_stat_card(slide, Inches(8.8), Inches(1.4), Inches(3.5), Inches(1.6),
                           pricing_src.title(), "Pricing Source", value_color=_DARK_BG)

        items = costs.get("items", [])
        if items:
            tbl_data = [["Service", "SKU", "Monthly Cost"]]
            for item in items:
                tbl_data.append([
                    item.get("service", ""),
                    item.get("sku", ""),
                    f"${item.get('monthly', 0):,.0f}",
                ])
            rows, cols = len(tbl_data), 3
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.45 * rows))
            tbl = table_shape.table
            tbl.columns[0].width = Inches(5)
            tbl.columns[1].width = Inches(3)
            tbl.columns[2].width = Inches(3.5)
            for r_idx, row_data in enumerate(tbl_data):
                for c_idx, val in enumerate(row_data):
                    cell = tbl.cell(r_idx, c_idx)
                    cell.text = str(val)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                        paragraph.font.name = "Calibri"
                        if r_idx == 0:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = _rgb("FFFFFF")
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _rgb(_TABLE_HDR)
                        else:
                            paragraph.font.color.rgb = _rgb(_TEXT)

    # ── Slide 7: Business Value ──────────────────────────────────────────
    bv = slide_data.get("businessValue")
    if bv:
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, _LIGHT_BG)
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                     "Business Value", font_size=30, bold=True, color=_DARK_BG)

        summary = bv.get("summary", "")
        if summary:
            _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(1.2),
                         summary, font_size=15, color=_TEXT)

        drivers = bv.get("drivers", [])
        if drivers:
            driver_strs = [f"{d.get('name', '')}: {d.get('impact', '')}" for d in drivers]
            _add_bullets(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(3.5), driver_strs)

        confidence = bv.get("confidence", "")
        if confidence:
            _add_textbox(slide, Inches(0.8), Inches(6.3), Inches(6), Inches(0.4),
                         f"Confidence Level: {confidence}", font_size=11, color=_SUBTEXT)

    # ── Slide 8: ROI ─────────────────────────────────────────────────────
    roi = slide_data.get("roi")
    if roi:
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, _LIGHT_BG)
        _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                     "Return on Investment", font_size=30, bold=True, color=_DARK_BG)

        roi_pct = roi.get("percent")
        payback = roi.get("paybackMonths")
        annual_val = roi.get("annualValue")
        annual_cost = roi.get("annualCost")

        col = 0
        card_w = Inches(2.8)
        gap = Inches(0.3)
        start_x = Inches(0.8)
        if roi_pct is not None:
            _add_stat_card(slide, start_x + col * (card_w + gap), Inches(1.5), card_w, Inches(1.8),
                           f"{roi_pct}%", "ROI")
            col += 1
        if payback is not None:
            _add_stat_card(slide, start_x + col * (card_w + gap), Inches(1.5), card_w, Inches(1.8),
                           f"{payback} mo", "Payback Period")
            col += 1
        if annual_val is not None:
            _add_stat_card(slide, start_x + col * (card_w + gap), Inches(1.5), card_w, Inches(1.8),
                           f"${annual_val:,.0f}", "Annual Value")
            col += 1
        if annual_cost is not None:
            _add_stat_card(slide, start_x + col * (card_w + gap), Inches(1.5), card_w, Inches(1.8),
                           f"${annual_cost:,.0f}", "Annual Cost")

    # ── Slide 9: Next Steps ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, _LIGHT_BG)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Next Steps", font_size=30, bold=True, color=_DARK_BG)
    next_steps = [
        "Review architecture and service selections with stakeholders",
        "Validate cost estimates against organizational procurement",
        "Conduct proof-of-concept for key components",
        "Develop detailed implementation timeline",
        "Engage Microsoft account team for enterprise agreements",
    ]
    _add_bullets(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(4), next_steps, font_size=15)

    # ── Slide 10: Closing ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, _DARK_BG)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), _ACCENT)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                 "Thank You", font_size=44, bold=True, color="FFFFFF", alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.8),
                 customer, font_size=24, color=_ACCENT, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "Generated by OneStopAgent", font_size=12, color=_SUBTEXT, alignment=PP_ALIGN.CENTER)

    prs.save(output_path)
    logger.info("Generated PPTX (python-pptx fallback): %s", output_path)
    return output_path
